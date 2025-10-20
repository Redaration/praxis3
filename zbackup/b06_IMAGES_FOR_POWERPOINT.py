#!/usr/bin/env python3
"""
PowerPoint Image Generator

This script reads a PowerPoint presentation and generates simple placeholder images for each slide,
then inserts these images into the slides.

Usage:
    To run as a script:
        python b06_IMAGES_FOR_POWERPOINT.py

    To use as a module:
        from b06_IMAGES_FOR_POWERPOINT import main
        main()

Requirements:
    - python-pptx: pip install python-pptx
    - pillow: pip install pillow
"""

import os
import sys
import time
import textwrap
from pptx import Presentation
from pptx.util import Inches
import io
from PIL import Image, ImageDraw, ImageFont
from random import choice

# Define constants
INPUT_PPTX = "course_presentation.pptx"
OUTPUT_PPTX = "course_presentation_with_images.pptx"
IMAGE_DIR = "slide_images"
MAX_SLIDES = 10  # Limit to processing only the first 10 slides

# Image settings - must be multiples of 64 for Runware API
IMAGE_WIDTH = 768    # Previously 800 - adjusted to be multiple of 64
IMAGE_HEIGHT = 576   # Previously 600 - adjusted to be multiple of 64
IMAGE_FORMAT = "png"

# Category colors and icons for the placeholder images
CATEGORY_STYLES = {
    "security": {"bg_color": "#2C3E50", "text_color": "#ECF0F1", "icon": "🔒"},
    "network": {"bg_color": "#3498DB", "text_color": "#FFFFFF", "icon": "🌐"},
    "cloud": {"bg_color": "#9B59B6", "text_color": "#FFFFFF", "icon": "☁️"},
    "technology": {"bg_color": "#2ECC71", "text_color": "#FFFFFF", "icon": "💻"},
    "corporate": {"bg_color": "#34495E", "text_color": "#FFFFFF", "icon": "🏢"},
    "default": {"bg_color": "#7F8C8D", "text_color": "#FFFFFF", "icon": "📊"}
}

def setup():
    """Set up the environment for image generation."""
    print("Setting up environment...")
    
    # Create directory for slide images if it doesn't exist
    if not os.path.exists(IMAGE_DIR):
        os.makedirs(IMAGE_DIR)
    
    # Check if PowerPoint file exists
    if not os.path.exists(INPUT_PPTX):
        print(f"Error: PowerPoint file '{INPUT_PPTX}' not found.")
        print("Please run 05_CREATE_POWERPOINT.py first to generate the presentation.")
        sys.exit(1)
        
    print(f"Will process only the first {MAX_SLIDES} slides for testing.")
    print("Using locally generated placeholder images.")



def get_image_category(slide_title, slide_content):
    """
    Determine the most appropriate image category based on slide content.
    
    Args:
        slide_title (str): The title of the slide
        slide_content (list): List of bullet points or content on the slide
        
    Returns:
        str: Category name for the image
    """
    # Convert everything to lowercase for easier keyword matching
    title_lower = slide_title.lower()
    content_text = ""
    if slide_content:
        content_text = " ".join(slide_content).lower()
    combined_text = title_lower + " " + content_text
    
    # Check for category keywords
    if any(word in combined_text for word in ["security", "protect", "threat", "risk", "compliance", "umbrella", "firewall"]):
        return "security"
    elif any(word in combined_text for word in ["network", "sd-wan", "routing", "connectivity", "vpn", "anyconnect"]):
        return "network"
    elif any(word in combined_text for word in ["cloud", "saas", "aws", "azure", "iaas", "paas", "sase"]):
        return "cloud"
    elif any(word in combined_text for word in ["technology", "digital", "innovation", "technical", "device", "hardware", "software"]):
        return "technology"
    elif any(word in combined_text for word in ["business", "management", "strategy", "corporate", "enterprise", "organization", "governance"]):
        return "corporate"
    else:
        return "default"

def create_placeholder_image(title, category, width, height):
    """
    Create a placeholder image with title text and category styling.
    
    Args:
        title (str): Title text to show on the image
        category (str): Category for styling (security, network, etc.)
        width (int): Image width in pixels
        height (int): Image height in pixels
        
    Returns:
        PIL.Image: The generated image
    """
    # Get style for the selected category
    style = CATEGORY_STYLES.get(category, CATEGORY_STYLES["default"])
    bg_color = style["bg_color"]
    text_color = style["text_color"]
    icon = style["icon"]
    
    # Create a new image with the specified color
    img = Image.new('RGB', (width, height), bg_color)
    draw = ImageDraw.Draw(img)
    
    # Try to load a font, falling back to default if not available
    try:
        title_font = ImageFont.truetype("Arial.ttf", 40)
        icon_font = ImageFont.truetype("Arial Unicode.ttf", 80)
    except IOError:
        title_font = ImageFont.load_default()
        icon_font = ImageFont.load_default()
    
    # Draw a category icon at the top
    icon_position = (width//2, height//4)
    draw.text(icon_position, icon, fill=text_color, font=icon_font, anchor="mm")
    
    # Prepare the title text - wrap it to multiple lines if needed
    max_chars_per_line = 30
    wrapped_title = textwrap.fill(title, max_chars_per_line)
    
    # Draw the title text
    text_position = (width//2, height//2)
    draw.text(text_position, wrapped_title, fill=text_color, font=title_font, align="center", anchor="mm")
    
    # Draw a subtle border
    border_width = 10
    for i in range(border_width):
        draw.rectangle(
            [(i, i), (width-i-1, height-i-1)],
            outline=text_color,
            width=1
        )
    
    return img

def generate_image_for_slide(slide_title, slide_content, output_path):
    """
    Generate an image based on slide title and content using Runware API.
    
    Args:
        slide_title (str): The title of the slide
        slide_content (list): List of bullet points or content on the slide
        output_path (str): Path where the generated image will be saved
    
    Returns:
        bool: True if successful, False otherwise
    """
    try:
        # Import the image generator module
        from runware_image_generator import generate_image
        
        print(f"Generating AI image for: {slide_title}")
        
        # Create a prompt based on slide content
        prompt = f"Create a professional, business-appropriate image that represents: {slide_title}"
        
        if slide_content:
            # Add more context from the slide content (limited to avoid making prompt too long)
            context = ". ".join(slide_content[:3])
            prompt += f". The image should relate to these concepts: {context}"
        
        # Add style guidance
        category = get_image_category(slide_title, slide_content)
        prompt += f". Style should be professional, clean design, suitable for a corporate presentation about {category}."
        
        print(f"Using prompt: {prompt}")
        
        # Generate the image using Runware
        downloaded_images = generate_image(
            prompt=prompt,
            output_path=output_path,
            width=IMAGE_WIDTH,
            height=IMAGE_HEIGHT,
            negative_prompt="text, watermarks, logos, inappropriate content, low quality"
        )
        
        # Check if image was successfully generated and downloaded
        if downloaded_images:
            print(f"AI image saved to {downloaded_images[0]}")
            return True
        else:
            # Fall back to local placeholder if API fails
            print("API call failed, falling back to local placeholder image")
            
            # Create a placeholder image
            image = create_placeholder_image(
                title=slide_title, 
                category=category,
                width=IMAGE_WIDTH,
                height=IMAGE_HEIGHT
            )
            
            # Save the image to a file
            image.save(output_path)
            print(f"Placeholder image saved to {output_path}")
            return True
    
    except Exception as e:
        print(f"Error generating image: {str(e)}")
        print("Falling back to local placeholder image")
        
        try:
            # Create a placeholder image as fallback
            category = get_image_category(slide_title, slide_content)
            image = create_placeholder_image(
                title=slide_title, 
                category=category,
                width=IMAGE_WIDTH,
                height=IMAGE_HEIGHT
            )
            
            # Save the image to a file
            image.save(output_path)
            print(f"Fallback image saved to {output_path}")
            return True
        except Exception as inner_e:
            print(f"Error creating fallback image: {str(inner_e)}")
            return False

def extract_slide_content(slide):
    """
    Extract title and content from a slide.
    
    Args:
        slide: A slide object from python-pptx
        
    Returns:
        tuple: (title_text, content_list)
    """
    title_text = ""
    content_list = []
    
    # Extract title text
    if slide.shapes.title:
        title_text = slide.shapes.title.text
    
    # Extract bullet points and other text content
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text and shape != slide.shapes.title:
            # Extract text from paragraphs
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text.strip():
                    content_list.append(paragraph.text.strip())
    
    return title_text, content_list

def add_image_to_slide(slide, image_path):
    """
    Add an image to a slide.
    
    Args:
        slide: A slide object from python-pptx
        image_path (str): Path to the image file
        
    Returns:
        bool: True if successful, False otherwise
    """
    try:
        # Calculate placement - position the image in the right half of the slide
        left = Inches(5.5)  # Position from the left edge (5.5 inches from left)
        top = Inches(1.5)   # Position from the top edge
        width = Inches(4)   # Width of the image
        
        # Add the image to the slide
        slide.shapes.add_picture(image_path, left, top, width=width)
        return True
    except Exception as e:
        print(f"Error adding image to slide: {str(e)}")
        return False

def process_presentation():
    """Process the presentation, generating and adding images for each slide."""
    try:
        # Load the presentation
        prs = Presentation(INPUT_PPTX)
        
        print(f"Processing presentation with {len(prs.slides)} slides (limiting to first {MAX_SLIDES})...")
        
        # Process each slide up to MAX_SLIDES limit
        slides_processed = 0
        for i, slide in enumerate(prs.slides):
            # Skip the title slide (first slide)
            if i == 0:
                print("Skipping title slide...")
                continue
            
            # Apply slide limit
            if slides_processed >= MAX_SLIDES:
                print(f"Reached max slides limit ({MAX_SLIDES}). Stopping image generation.")
                break
                
            slides_processed += 1
            
            # Extract slide content
            title, content = extract_slide_content(slide)
            if not title:
                print(f"Skipping slide {i+1} (no title)")
                slides_processed -= 1  # Don't count this one
                continue
            
            # Define image file path
            image_filename = f"slide_{i+1:03d}_{title.replace(' ', '_')[:30]}.{IMAGE_FORMAT}"
            image_path = os.path.join(IMAGE_DIR, image_filename)
            
            # Generate image if it doesn't exist
            if not os.path.exists(image_path):
                success = generate_image_for_slide(title, content, image_path)
                if not success:
                    print(f"Failed to generate image for slide {i+1}")
                    continue
                
                # Add a short delay to avoid rate limiting
                time.sleep(2)
            else:
                print(f"Using existing image for slide {i+1}: {image_path}")
            
            # Add the image to the slide
            add_image_to_slide(slide, image_path)
        
        # Save the modified presentation
        prs.save(OUTPUT_PPTX)
        print(f"\nPresentation with images saved as '{OUTPUT_PPTX}'")
        print(f"Processed {slides_processed} slides out of {len(prs.slides)-1} total content slides")
        
    except Exception as e:
        print(f"Error processing presentation: {str(e)}")
        raise

def main():
    """Main function to run the image generation and insertion process."""
    print("=" * 70)
    print("PowerPoint Image Generator".center(70))
    print("=" * 70)
    
    # Set up the environment
    setup()
    
    # Process the presentation
    process_presentation()
    
    print("\nDone!")
    print("=" * 70)

if __name__ == "__main__":
    main()
